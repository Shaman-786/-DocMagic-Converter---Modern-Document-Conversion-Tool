import os
from pptx import Presentation
from pdf2image import convert_from_path
from utils.file_handling import get_output_path


def convert(input_path, output_folder):
    output_path = get_output_path(input_path, output_folder, ".pptx")

    # Convert PDF to images
    images = convert_from_path(input_path)

    # Create PowerPoint
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Blank slide

    for image in images:
        # Save temp image
        temp_img_path = os.path.join(output_folder, "temp.jpg")
        image.save(temp_img_path, "JPEG")

        # Add slide with image
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            temp_img_path,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )

        # Remove temp image
        os.remove(temp_img_path)

    prs.save(output_path)
    return output_path